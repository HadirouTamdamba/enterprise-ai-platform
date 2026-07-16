"""Document parsers for enterprise formats (F-21). Each returns (page, text) tuples."""

import csv
import io
import json
from pathlib import Path

from app.core.exceptions import ValidationError

PageText = tuple[int | None, str]

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".xlsx", ".md", ".markdown", ".txt",
    ".csv", ".json", ".html", ".htm", ".png", ".jpg", ".jpeg",
}


def parse_document(path: Path, content_type: str = "") -> list[PageText]:
    """Dispatch to a format-specific parser based on file extension."""
    suffix = path.suffix.lower()
    parser = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".pptx": _parse_pptx,
        ".xlsx": _parse_xlsx,
        ".csv": _parse_csv,
        ".json": _parse_json,
        ".html": _parse_html,
        ".htm": _parse_html,
        ".md": _parse_text,
        ".markdown": _parse_text,
        ".txt": _parse_text,
        ".png": _parse_image_ocr,
        ".jpg": _parse_image_ocr,
        ".jpeg": _parse_image_ocr,
    }.get(suffix)
    if parser is None:
        raise ValidationError(f"Unsupported document format: {suffix}")
    return parser(path)


def _parse_pdf(path: Path) -> list[PageText]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages: list[PageText] = []
    for number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if not text.strip():  # scanned page → OCR fallback if available
            text = _ocr_pdf_page(path, number)
        pages.append((number, text))
    return pages


def _parse_docx(path: Path) -> list[PageText]:
    import docx

    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return [(None, "\n".join(parts))]


def _parse_pptx(path: Path) -> list[PageText]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    pages: list[PageText] = []
    for number, slide in enumerate(presentation.slides, start=1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        pages.append((number, "\n".join(texts)))
    return pages


def _parse_xlsx(path: Path) -> list[PageText]:
    import openpyxl

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    pages: list[PageText] = []
    for index, sheet in enumerate(workbook.worksheets, start=1):
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        pages.append((index, f"Sheet: {sheet.title}\n" + "\n".join(rows)))
    return pages


def _parse_csv(path: Path) -> list[PageText]:
    with path.open(newline="", encoding="utf-8", errors="replace") as handle:
        rows = [" | ".join(row) for row in csv.reader(handle)]
    return [(None, "\n".join(rows))]


def _parse_json(path: Path) -> list[PageText]:
    data = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    return [(None, json.dumps(data, indent=2, ensure_ascii=False))]


def _parse_html(path: Path) -> list[PageText]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="replace"), "html.parser")
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    return [(None, soup.get_text(separator="\n", strip=True))]


def _parse_text(path: Path) -> list[PageText]:
    return [(None, path.read_text(encoding="utf-8", errors="replace"))]


def _parse_image_ocr(path: Path) -> list[PageText]:
    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise ValidationError(
            "OCR support requires the 'ocr' extra: pip install '.[ocr]'"
        ) from exc
    return [(None, pytesseract.image_to_string(Image.open(path)))]


def _ocr_pdf_page(path: Path, page_number: int) -> str:
    """Best-effort OCR for scanned PDF pages; returns empty text when OCR is not installed."""
    try:
        import pytesseract  # noqa: F401
    except ImportError:
        return ""
    # Rendering PDFs to images requires poppler/pdf2image; documented in the ocr extra.
    try:
        from pdf2image import convert_from_path

        images = convert_from_path(str(path), first_page=page_number, last_page=page_number)
        if images:
            import pytesseract

            return pytesseract.image_to_string(images[0])
    except Exception:
        return ""
    return ""


def extract_metadata(path: Path, content_type: str) -> dict:
    stat = path.stat()
    return {
        "extension": path.suffix.lower(),
        "content_type": content_type,
        "size_bytes": stat.st_size,
    }
